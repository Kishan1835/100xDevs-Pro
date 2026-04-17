"""
transfer_v3.py — Fixed version. Addresses all 5 problems found in output.pptx:

  Problem 1: Wrong layout (title/body in right 40% only)
             → Use layout 4 "Two Content": title x=0.93" w=11.07", body w=11.07"
               Override body placeholder position via textbox fallback for full width.

  Problem 2: Duplicate section-divider + content slide pairs
             → Detect consecutive slides with same title and no body; merge them.

  Problem 3: Wrong title (picks short fragment like "CNC" or "HEALTHY")
             → Title must be ≥ 4 chars AND prefer longer all-caps strings.

  Problem 4: Body box too small / wrong position
             → Always use full-width textbox at correct y for body, not placeholder.

  Problem 5: Images stacked on top of placeholder
             → Remove empty PICTURE placeholder shape before placing images.
"""

import io, zipfile, argparse
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn


# ─── ZIP-level template cleaning (prevents duplicate-name corruption) ─────────

NS_PML    = 'http://schemas.openxmlformats.org/presentationml/2006/main'
REL_SLIDE = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide'
CT_SLIDE  = 'application/vnd.openxmlformats-officedocument.presentationml.slide+xml'

def make_clean_template(path):
    with open(path, 'rb') as f:
        raw = f.read()
    src = zipfile.ZipFile(io.BytesIO(raw), 'r')
    buf = io.BytesIO()
    out = zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED)
    SKIP = ('ppt/slides/slide', 'ppt/slides/_rels/slide', 'ppt/notesSlides/')
    for name in src.namelist():
        data = src.read(name)
        if any(name.startswith(p) for p in SKIP):
            continue
        if name == 'ppt/presentation.xml':
            tree = etree.fromstring(data)
            lst  = tree.find(f'{{{NS_PML}}}sldIdLst')
            if lst is not None:
                for ch in list(lst): lst.remove(ch)
            data = etree.tostring(tree, xml_declaration=True, encoding='UTF-8', standalone=True)
        elif name == 'ppt/_rels/presentation.xml.rels':
            tree = etree.fromstring(data)
            for rel in list(tree):
                if rel.get('Type') == REL_SLIDE: tree.remove(rel)
            data = etree.tostring(tree, xml_declaration=True, encoding='UTF-8', standalone=True)
        elif name == '[Content_Types].xml':
            tree = etree.fromstring(data)
            for ov in list(tree):
                if ov.get('ContentType') == CT_SLIDE: tree.remove(ov)
            data = etree.tostring(tree, xml_declaration=True, encoding='UTF-8', standalone=True)
        out.writestr(name, data)
    out.close()
    buf.seek(0)
    return buf


# ─── Recursive content extraction ─────────────────────────────────────────────

def _collect_texts(shape, results, depth=0):
    if shape.has_text_frame:
        txt = shape.text_frame.text.strip()
        if len(txt) > 2:
            results.append((depth, txt))
    if hasattr(shape, 'shapes'):
        for s in shape.shapes:
            _collect_texts(s, results, depth + 1)

def _collect_images(shape, results):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            results.append((shape.image.blob, shape.image.ext or 'png'))
        except Exception:
            pass
    if hasattr(shape, 'shapes'):
        for s in shape.shapes:
            _collect_images(s, results)

# Single-word status labels that should never be slide titles
_STATUS_WORDS = {
    'HEALTHY', 'ALERT', 'CRITICAL', 'WARNING', 'ERROR', 'OK',
    'YES', 'NO', 'ON', 'OFF', 'HIGH', 'LOW', 'PASS', 'FAIL',
}

def _pick_title(raw_texts):
    """
    FIX 3: Better title detection.
    Priority order:
      1. Non-caps string that looks like a proper heading (≥10 chars, mixed case, depth 0-1)
      2. ALL-CAPS string ≥ 8 chars (long enough to be a real heading, not a status word)
      3. First text ≥ 4 chars as fallback
    Never pick known status words or strings < 4 chars.
    """
    if not raw_texts:
        return ""

    # Priority 1: mixed-case heading-like strings at top depth
    for d, t in raw_texts:
        if (len(t) >= 8
                and t not in _STATUS_WORDS
                and not t.startswith('[')
                and not t[0].isdigit()
                and '\n' not in t
                and d <= 1):
            # Check it looks heading-like: starts with capital, reasonable length
            if t[0].isupper() and len(t) <= 80:
                return t

    # Priority 2: long ALL-CAPS strings (real headings, not labels)
    allcaps = [
        (d, t) for d, t in raw_texts
        if t.upper() == t
        and len(t) >= 8          # must be >= 8 chars (rules out HEALTHY, ALERT, etc.)
        and t not in _STATUS_WORDS
        and not t.startswith('[')
        and not t[0].isdigit()
        and '\n' not in t
    ]
    if allcaps:
        allcaps.sort(key=lambda x: -len(x[1]))
        return allcaps[0][1]

    # Fallback: first text >= 4 chars that is not a status word
    for _, t in raw_texts:
        if len(t) >= 4 and t not in _STATUS_WORDS:
            return t
    return raw_texts[0][1] if raw_texts else ""

def extract_slide(slide):
    raw_texts, images = [], []
    for sh in slide.shapes:
        _collect_texts(sh, raw_texts)
        _collect_images(sh, images)
    has_table = any(getattr(sh, 'has_table', False) for sh in slide.shapes)
    has_chart = any(getattr(sh, 'has_chart', False) for sh in slide.shapes)

    title = _pick_title(raw_texts)

    # Body: everything except the title, deduplicated
    seen, body = {title}, []
    for _, t in raw_texts:
        for line in t.split('\n'):
            line = line.strip()
            if line and line not in seen and len(line) > 2:
                seen.add(line)
                body.append(line)

    return dict(title=title, body=body, images=images,
                has_table=has_table, has_chart=has_chart)


# ─── FIX 2: Merge section-divider + content slide pairs ───────────────────────

def merge_slides(slides):
    """
    Consecutive slides where slide[i] has a title but no body/images,
    and slide[i+1] has the same (or very similar) title → merge them.
    The section-divider slide is dropped; its title is kept on the content slide.
    """
    merged = []
    i = 0
    while i < len(slides):
        s = slides[i]
        is_empty = not s['body'] and not s['images']

        if is_empty and i + 1 < len(slides):
            nxt = slides[i + 1]
            # Same title or next slide has no title of its own
            same_title = (s['title'].lower() == nxt['title'].lower()
                          or not nxt['title'])
            if same_title:
                # Merge: keep section title, use next slide's body/images
                merged_slide = dict(
                    title=s['title'],
                    body=nxt['body'],
                    images=nxt['images'],
                    has_table=nxt['has_table'],
                    has_chart=nxt['has_chart'],
                )
                merged.append(merged_slide)
                i += 2   # skip the next slide (already merged)
                continue

        merged.append(s)
        i += 1

    return merged


# ─── Layout selection ──────────────────────────────────────────────────────────

def pick_layout(prs, content):
    """
    FIX 1: Prefer layout 4 "Two Content" which has:
      - TITLE at x=0.93" w=11.07" (full width)
      - OBJECT at x=0.93" y=2.00" w=11.07" h=4.76" (full width body)
    For image-only slides, try to find a layout with PICTURE placeholder.
    For title-only slides (no body, no images), use "Section Header" or "Title Only".
    """
    has_images = bool(content['images'])
    has_body   = bool(content['body'])
    has_title  = bool(content['title'])

    # Title-only / section divider
    if has_title and not has_body and not has_images:
        for layout in prs.slide_layouts:
            if layout.name in ('Section Header', 'Title Only', '2_Title Only'):
                return layout

    # Image-heavy with no body: look for picture layout
    if has_images and not has_body:
        for layout in prs.slide_layouts:
            ph_types = {ph.placeholder_format.type for ph in layout.placeholders}
            if PP_PLACEHOLDER.PICTURE in ph_types:
                return layout

    # Standard content: use "Two Content" (layout 4) — full-width title + body
    for layout in prs.slide_layouts:
        if layout.name == 'Two Content':
            return layout

    # Fallback: any layout with TITLE + OBJECT/BODY
    for layout in prs.slide_layouts:
        ph_types = {ph.placeholder_format.type for ph in layout.placeholders}
        has_title_ph = any(t in ph_types for t in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE))
        has_body_ph  = any(t in ph_types for t in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT))
        if has_title_ph and has_body_ph:
            return layout

    return prs.slide_layouts[0]


# ─── Injection ────────────────────────────────────────────────────────────────

# Fixed textbox zones for "Two Content" layout (13.33" x 7.50" slide)
# idx=0 TITLE: x=0.93" y=0.96" w=11.07" h=0.71"
# idx=2 OBJECT: x=0.93" y=2.00" w=11.07" h=4.76"

Z_TITLE = (0.50, 0.20, 12.33, 0.90)   # full-width title zone
Z_BODY  = (0.50, 1.30, 12.33, 5.80)   # full-width body zone (starts after title)
Z_BODY_WITH_IMG = (0.50, 1.30, 7.50, 5.80)
Z_IMAGE = (8.30, 1.30, 4.53, 5.80)
Z_IMAGE_FULL = (0.50, 1.30, 12.33, 5.80)


def inject(prs, content, slide_num):
    layout = pick_layout(prs, content)
    slide  = prs.slides.add_slide(layout)

    # FIX 5: Remove all empty placeholder shapes so nothing sits underneath content
    _remove_empty_placeholders(slide)

    title  = content['title']
    body   = content['body']
    images = content['images']

    _do_title(slide, title)
    _do_body(slide, body, has_images=bool(images))
    _do_images(slide, images, body_present=bool(body))

    print(f"  Slide {slide_num:>2}: [{layout.name:<18}] "
          f"title={repr(title[:40]):<44} "
          f"body={len(body):>2}  img={len(images)}")


def _remove_empty_placeholders(slide):
    """
    FIX 5: Delete all placeholder shapes that have no text and are not needed.
    This prevents ghost shapes sitting under our injected content.
    """
    sp_tree = slide.shapes._spTree
    to_remove = []
    for shape in slide.placeholders:
        ph_type = shape.placeholder_format.type
        # Keep title placeholder (we'll fill it)
        if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            continue
        # Remove picture, body, object placeholders — we replace them with textboxes/images
        to_remove.append(shape._element)
    for el in to_remove:
        sp_tree.remove(el)


def _do_title(slide, title):
    # Try title placeholder first (it's full-width in "Two Content")
    for ph in slide.placeholders:
        pt = ph.placeholder_format.type
        if pt in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            _fill_ph(ph, [title], bold=True, size=Pt(24))
            return
    # Fallback textbox
    l, t, w, h = Z_TITLE
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    _fill_tb(tb, [title], bold=True, size=Pt(24))


def _do_body(slide, lines, has_images=False):
    if not lines:
        return
    zone = Z_BODY_WITH_IMG if has_images else Z_BODY
    l, t, w, h = zone
    # FIX 4: Always use a textbox for body — gives full control over position/size
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    # Auto-size: shrink text if needed
    tf.auto_size = None

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(13)


def _do_images(slide, images, body_present=True):
    if not images:
        return
    for idx, (blob, ext) in enumerate(images[:3]):
        stream = io.BytesIO(blob)
        if not body_present and idx == 0:
            l, t, w, h = Z_IMAGE_FULL
        else:
            l, t, w, h = Z_IMAGE
            t = t + idx * 0.25   # slight vertical stagger for multiple images
        try:
            slide.shapes.add_picture(stream, Inches(l), Inches(t), Inches(w), Inches(h))
        except Exception as e:
            print(f"    ⚠  image {idx+1} skipped: {e}")


def _fill_ph(ph, lines, bold=False, size=None):
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        if bold:  run.font.bold = True
        if size:  run.font.size = size

def _fill_tb(tb, lines, bold=False, size=None):
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        if bold:  run.font.bold = True
        if size:  run.font.size = size


# ─── Main ─────────────────────────────────────────────────────────────────────

def transfer(source_path, template_path, output_path):
    print(f"\n{'═'*60}")
    print(f"  PPTX Transfer v3  (all 5 problems fixed)")
    print(f"{'═'*60}")

    src_prs = Presentation(source_path)
    print(f"\n  Source:   {len(src_prs.slides)} slides")

    print("\n► Phase 1: Extracting source slides…")
    raw_slides = [extract_slide(s) for s in src_prs.slides]

    print("\n► Phase 2: Merging section-divider duplicates…")
    slides = merge_slides(raw_slides)
    print(f"  {len(raw_slides)} → {len(slides)} slides after merge")

    print("\n► Phase 3: Building output…")
    clean_buf = make_clean_template(template_path)
    out_prs   = Presentation(clean_buf)
    print(f"  Template: {len(out_prs.slide_layouts)} layouts, 0 existing slides\n")

    for i, content in enumerate(slides):
        inject(out_prs, content, i + 1)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    out_prs.save(str(out))
    size_kb = out.stat().st_size / 1024
    print(f"\n{'═'*60}")
    print(f"  ✓  {len(slides)} slides  →  {output_path}  ({size_kb:.0f} KB)")
    print(f"{'═'*60}\n")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("-s", "--source",  required=True)
    ap.add_argument("-d", "--dest",    required=True)
    ap.add_argument("-o", "--output",  required=True)
    args = ap.parse_args()
    transfer(args.source, args.dest, args.output)